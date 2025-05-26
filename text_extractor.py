# text_extractor.py
import os
import pdfplumber
import docx
import pptx

def extract_text_from_pdf(file_path: str) -> str | None:
    """Extracts text content from a PDF file."""
    try:
        with pdfplumber.open(file_path) as pdf:
            text_content = "".join([page.extract_text() + "\n" for page in pdf.pages if page.extract_text()])
            return text_content
    except Exception as e:
        print(f"Extractor: Error reading PDF {file_path}: {e}")
        return None

def extract_text_from_txt_or_code(file_path: str) -> str | None:
    """Extracts text content from a .txt, .md, .cpp, .h, .js, .py file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Extractor: Error reading text/code {file_path}: {e}")
        return None

def extract_text_from_docx(file_path: str) -> str | None:
    """Extracts text content from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Extractor: Error reading DOCX {file_path}: {e}")
        return None

def extract_text_from_pptx(file_path: str) -> str | None:
    """Extracts text content from a PPTX file."""
    try:
        prs = pptx.Presentation(file_path)
        return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except Exception as e:
        print(f"Extractor: Error reading PPTX {file_path}: {e}")
        return None

def extract_text(file_path: str) -> str | None:
    """Extracts text from various file types based on extension."""
    filename = os.path.basename(file_path)
    file_ext = filename.lower().split('.')[-1] if '.' in filename else ''

    if file_ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_ext in ["txt", "md", "cpp", "h", "js", "py", "gitignore", "log"]:
        return extract_text_from_txt_or_code(file_path)
    elif file_ext == "docx":
        return extract_text_from_docx(file_path)
    elif file_ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file_path)
    else:
        # print(f"Extractor: Skipping unsupported file type: {filename}")
        return None