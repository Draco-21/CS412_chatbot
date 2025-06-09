import os
import re
import pdfplumber  # For reading PDF files
import docx        # For reading DOCX files
import pptx        # For reading PPTX files
import PyPDF2
import json
import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Set
from dataclasses import dataclass
from functools import lru_cache
import concurrent.futures
import hashlib
from threading import Lock
import nltk
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.stem import WordNetLemmatizer
from nltk.corpus import stopwords
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer
import numpy as np
from collections import defaultdict
from chatbotcodes.models import ResourceMatch

RESOURCES_BASE_FOLDER = "cleaned_resources"
YEAR_FOLDER_MAP = {
    "Year 1 Certificate": "Y1",
    "Year 2 Diploma": "Y2",
    "Year 3 Degree": "Y3",
    "Year 4 Postgraduate Diploma": "Y4"
}

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('chatbot.log')
    ]
)
logger = logging.getLogger(__name__)

# Initialize NLTK
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
    nltk.data.find('wordnet')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')
    nltk.download('wordnet')

# --- Text Extraction Functions (No changes needed here) ---

def extract_text_from_pdf(file_path: str) -> str | None:
    """Extracts text content from a PDF file."""
    try:
        with pdfplumber.open(file_path) as pdf:
            text_content = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text + "\n"
            return text_content.lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading PDF file {file_path}: {e}")
        return None

def extract_text_from_txt_or_code(file_path: str) -> str | None:
    """Extracts text content from a .txt, .md, .cpp, .h, .js, .py file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading text/code file {file_path}: {e}")
        return None

def extract_text_from_docx(file_path: str) -> str | None:
    """Extracts text content from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text).lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading DOCX file {file_path}: {e}")
        return None

def extract_text_from_pptx(file_path: str) -> str | None:
    """Extracts text content from a PPTX file."""
    try:
        prs = pptx.Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text).lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading PPTX file {file_path}: {e}")
        return None

# --- NEW Helper Function to Search a Single Folder ---

def search_single_folder(folder_path: str, keywords: list[str]) -> str | None:
    """
    Searches recursively within a single folder path for files matching keywords.
    Returns the full content of the first matching file found.
    """
    if not os.path.isdir(folder_path):
        print(f"ResourceRetriever: Search path {folder_path} not found or not a directory.")
        return None

    try:
        for root, dirs, files in os.walk(folder_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                content = None
                file_ext = filename.lower().split('.')[-1] if '.' in filename else ''

                # Extract text based on file type
                if file_ext == "pdf":
                    content = extract_text_from_pdf(file_path)
                elif file_ext in ["txt", "md", "cpp", "h", "js", "py", "gitignore", "log"]:
                    content = extract_text_from_txt_or_code(file_path)
                elif file_ext == "docx":
                    content = extract_text_from_docx(file_path)
                elif file_ext in ["ppt", "pptx"]:
                    content = extract_text_from_pptx(file_path)

                if content:
                    # Check if ANY keyword is present
                    found_keywords_count = sum(1 for keyword in keywords if keyword in content)
                    if found_keywords_count > 0:
                        print(f"ResourceRetriever: Found {found_keywords_count} keyword(s) in '{file_path}'.")
                        return content # Return full content of the first match

    except Exception as e:
        print(f"ResourceRetriever: Error during search in {folder_path}: {e}")
        return None

    return None # No match found in this folder

# --- UPDATED Main Search Function ---

def find_answer_in_resources(user_message: str, year_level: str) -> str | None:
    """
    Searches for an answer within local resource files, checking the primary
    year level first, then expanding to other year levels.
    """
    print(f"ResourceRetriever: Main search for '{user_message}', primary year: {year_level}.")

    keywords = [word for word in re.split(r'\W+', user_message.lower()) if word]
    if not keywords:
        print("ResourceRetriever: No valid keywords in user message.")
        return None
    
    print(f"ResourceRetriever: Using keywords: {keywords}")

    # --- Determine Search Order ---
    primary_folder_name = YEAR_FOLDER_MAP.get(year_level)
    all_folder_names = list(YEAR_FOLDER_MAP.values())
    
    search_order_names = []
    
    # 1. Add primary folder first if it exists
    if primary_folder_name:
        search_order_names.append(primary_folder_name)
    else:
        print(f"ResourceRetriever: Warning - No folder mapping for {year_level}")

    # 2. Add all other folders (excluding the primary one if it was added)
    for name in all_folder_names:
        if name not in search_order_names:
            search_order_names.append(name)

    print(f"ResourceRetriever: Search order: {search_order_names}")

    # --- Execute Search in Order ---
    for folder_name in search_order_names:
        folder_path = os.path.join(RESOURCES_BASE_FOLDER, folder_name)
        print(f"ResourceRetriever: Now searching in: {folder_path}")
        
        result = search_single_folder(folder_path, keywords)
        
        if result:
            print(f"ResourceRetriever: Match found in {folder_path}.")
            return result # Return the first match we find

    # If we get here, no match was found in any folder
    print(f"ResourceRetriever: No match found in any folder for '{user_message}'.")
    return None

# --- Example Usage (remains the same for testing) ---
if __name__ == "__main__":
    # ... (Keep your example usage/dummy file creation if needed) ...
    print("Testing with 'C++ loops'")
    result = find_answer_in_resources("C++ loops", "Year 1 Certificate")
    if result:
        print("\n--- Search Result Found (Showing first 500 chars) ---")
        print(result[:500] + "...")
    else:
        print("\nNo local result found.")

@dataclass
class ResourceMatch:
    title: str
    language: str
    purpose: str
    code: str
    context: str
    source_file: str
    topics: List[str]
    framework_name: str = ""
    framework_version: str = ""
    difficulty_level: str = ""
    score: float = 0.0

class SmartResourceRetriever:
    def __init__(self, resources_dir=None):
        logger.info("Initializing SmartResourceRetriever")
        self.resources_dir = resources_dir or RESOURCES_BASE_FOLDER
        self._content_cache = {}  # File path -> content cache
        self._cache_lock = Lock()  # Thread-safe cache access
        self._keyword_cache = {}   # Query -> keywords cache
        self._embedding_cache = {}  # Text -> embedding cache
        
        # Initialize NLP tools
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words('english'))
        self.tfidf_vectorizer = TfidfVectorizer(
            tokenizer=self._preprocess_text,
            stop_words='english',
            ngram_range=(1, 2)
        )
        
        # Initialize sentence transformer model for semantic embeddings
        self.sentence_transformer = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Enhanced topic categories with hierarchical structure
        self.topic_categories = {
            'fundamentals': {
                'basic_concepts': ['syntax', 'variable', 'datatype', 'operator'],
                'control_flow': ['loop', 'condition', 'function', 'recursion'],
                'error_handling': ['exception', 'try', 'catch', 'debug']
            },
            'data_structures': {
                'linear': ['array', 'list', 'stack', 'queue'],
                'hierarchical': ['tree', 'heap', 'binary'],
                'graph': ['graph', 'network', 'vertex', 'edge'],
                'hash': ['hash', 'map', 'dictionary', 'set']
            },
            'algorithms': {
                'searching': ['search', 'binary search', 'linear search'],
                'sorting': ['sort', 'quicksort', 'mergesort'],
                'optimization': ['dynamic', 'greedy', 'backtrack'],
                'complexity': ['time complexity', 'space complexity', 'big o']
            },
            'web_dev': {
                'frontend': ['html', 'css', 'javascript', 'react'],
                'backend': ['api', 'server', 'database', 'rest'],
                'security': ['authentication', 'authorization', 'encryption']
            }
        }
        
        # Enhanced programming language detection
        self.programming_languages = {
            'python': {
                'core': ['python', 'def', 'class', 'import'],
                'frameworks': ['django', 'flask', 'fastapi'],
                'libraries': ['pandas', 'numpy', 'tensorflow']
            },
            'cpp': {
                'core': ['c++', 'cpp', 'iostream', 'vector'],
                'stl': ['std::', 'template', 'container'],
                'features': ['pointer', 'reference', 'memory']
            },
            'java': {
                'core': ['java', 'class', 'interface', 'package'],
                'frameworks': ['spring', 'hibernate', 'junit'],
                'features': ['inheritance', 'polymorphism', 'encapsulation']
            },
            'javascript': {
                'core': ['js', 'function', 'const', 'let'],
                'frontend': ['dom', 'event', 'react', 'vue'],
                'backend': ['node', 'express', 'npm']
            }
        }
        
        # Load and cache all resources on initialization
        self._load_all_resources()
        
        logger.info(f"SmartResourceRetriever initialized with {len(self._content_cache)} resources")

    def format_matches_for_display(self, matches: List[ResourceMatch]) -> str:
        """Format resource matches into a readable string format."""
        if not matches:
            return "No relevant matches found."
        
        formatted_output = []
        for i, match in enumerate(matches, 1):
            formatted_match = f"\n### Match {i}: {match.title}\n"
            formatted_match += f"Language: {match.language}\n"
            formatted_match += f"Purpose: {match.purpose}\n"
            
            if match.context:
                formatted_match += f"Context: {match.context}\n"
            
            formatted_match += f"\n```{match.language}\n{match.code}\n```\n"
            
            if match.topics:
                formatted_match += f"\nTopics: {', '.join(match.topics)}\n"
            
            if match.framework_name:
                formatted_match += f"Framework: {match.framework_name}"
                if match.framework_version:
                    formatted_match += f" (v{match.framework_version})"
                formatted_match += "\n"
            
            if match.difficulty_level:
                formatted_match += f"Difficulty: {match.difficulty_level}\n"
            
            formatted_match += f"Relevance Score: {match.score:.2f}\n"
            formatted_output.append(formatted_match)
        
        return "\n".join(formatted_output)

    def _load_all_resources(self):
        """Load all JSON resources into memory cache."""
        try:
            # Get absolute path to resources directory
            resources_root = Path(os.path.abspath(self.resources_dir))
            if not resources_root.exists():
                logger.error(f"Resources directory not found: {resources_root}")
                return

            for json_file in resources_root.rglob('*.json'):
                try:
                    with open(json_file, 'r', encoding='utf-8') as f:
                        content = json.load(f)
                        self._content_cache[str(json_file)] = content
                        logger.debug(f"Loaded resource: {json_file}")
                except Exception as e:
                    logger.error(f"Error loading {json_file}: {str(e)}")

        except Exception as e:
            logger.error(f"Error loading resources: {str(e)}")

    def _preprocess_text(self, text: str) -> List[str]:
        """Enhanced text preprocessing with better tokenization and normalization."""
        # Convert to lowercase and tokenize
        tokens = word_tokenize(text.lower())
        
        # Custom handling for common programming terms
        programming_terms = {
            'programming': 'program',
            'programs': 'program',
            'algorithmic': 'algorithm',
            'algorithms': 'algorithm'
        }
        
        # Remove stopwords and lemmatize with custom handling
        processed_tokens = []
        for token in tokens:
            if token not in self.stop_words and token.isalnum():
                # Check for custom mappings first
                if token in programming_terms:
                    processed_tokens.append(programming_terms[token])
                else:
                    lemmatized = self.lemmatizer.lemmatize(token)
                    processed_tokens.append(lemmatized)
        
        # Add bigrams for better context
        bigrams = [f"{processed_tokens[i]}_{processed_tokens[i+1]}" 
                  for i in range(len(processed_tokens)-1)]
        
        return processed_tokens + bigrams

    def _get_text_embedding(self, text: str) -> np.ndarray:
        """Get or compute embedding for text using sentence transformer."""
        if text in self._embedding_cache:
            return self._embedding_cache[text]
            
        embedding = self.sentence_transformer.encode(text, convert_to_tensor=True)
        self._embedding_cache[text] = embedding
        return embedding

    def _calculate_semantic_similarity(self, text1: str, text2: str) -> float:
        """Calculate semantic similarity between two texts using embeddings."""
        emb1 = self._get_text_embedding(text1)
        emb2 = self._get_text_embedding(text2)
        return float(cosine_similarity(emb1.reshape(1, -1), emb2.reshape(1, -1))[0][0])

    def _detect_topics_and_language(self, query: str) -> Dict[str, Dict]:
        """Enhanced topic and language detection with hierarchical matching."""
        query_tokens = set(self._preprocess_text(query))
        
        detected = {
            'topics': defaultdict(list),
            'languages': defaultdict(list)
        }
        
        # Detect topics with hierarchical matching and fuzzy matching
        for main_category, subcategories in self.topic_categories.items():
            for subcategory, keywords in subcategories.items():
                # Check for exact matches
                matches = [keyword for keyword in keywords 
                          if any(token in query_tokens for token in keyword.split())]
                
                # Check for partial matches (e.g., "algorithm" matches "algorithms")
                partial_matches = [keyword for keyword in keywords 
                                 if any(token.startswith(keyword) or keyword.startswith(token) 
                                      for token in query_tokens)]
                
                all_matches = list(set(matches + partial_matches))
                if all_matches:
                    detected['topics'][main_category].append({
                        'subcategory': subcategory,
                        'matches': all_matches
                    })
        
        # Enhanced language detection with context awareness
        for lang, categories in self.programming_languages.items():
            for category, keywords in categories.items():
                # Check for exact and partial matches
                matches = [keyword for keyword in keywords 
                          if any(token in query_tokens or
                                any(token.startswith(k) or k.startswith(token) 
                                    for k in keyword.split())
                                for token in query_tokens)]
                if matches:
                    detected['languages'][lang].append({
                        'category': category,
                        'matches': matches
                    })
        
        return detected

    def _calculate_relevance_score(self, content: Dict, query: str,
                                 detected_info: Dict[str, Dict]) -> float:
        """
        Calculate an enhanced relevance score using multiple factors:
        1. TF-IDF similarity
        2. Semantic similarity using sentence embeddings
        3. Topic and language relevance
        4. Code quality and structure
        """
        # Initialize score components
        tfidf_score = 0.0
        semantic_score = 0.0
        topic_score = 0.0
        language_score = 0.0
        
        # Calculate TF-IDF similarity
        if content.get('text'):
            try:
                tfidf_matrix = self.tfidf_vectorizer.fit_transform([query, content['text']])
                tfidf_score = float(cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0])
            except Exception as e:
                logger.warning(f"Error calculating TF-IDF score: {e}")
        
        # Calculate semantic similarity
        try:
            semantic_score = self._calculate_semantic_similarity(query, content.get('text', ''))
        except Exception as e:
            logger.warning(f"Error calculating semantic similarity: {e}")
        
        # Calculate topic relevance
        if detected_info['topics']:
            topic_matches = sum(len(subcats) for subcats in detected_info['topics'].values())
            topic_score = min(1.0, topic_matches / 3.0)  # Normalize to [0,1]
        
        # Calculate language relevance
        if detected_info['languages']:
            lang_matches = sum(len(cats) for cats in detected_info['languages'].values())
            language_score = min(1.0, lang_matches / 2.0)  # Normalize to [0,1]
        
        # Weighted combination of scores
        final_score = (
            0.3 * tfidf_score +
            0.3 * semantic_score +
            0.2 * topic_score +
            0.2 * language_score
        )
        
        return min(1.0, max(0.0, final_score))  # Ensure score is in [0,1]

    def search_resources(self, query: str, year_level: str, min_score: float = 0.2) -> List[ResourceMatch]:
        """
        Search through resources using enhanced semantic search and scoring.
        Returns list of ResourceMatch objects sorted by relevance score.
        """
        logger.info(f"Searching for: '{query}' (year level: {year_level})")
        
        # Preprocess query and detect topics/language
        query_tokens = self._preprocess_text(query)
        logger.debug(f"Query tokens: {query_tokens}")
        
        detected_info = self._detect_topics_and_language(query)
        logger.debug(f"Detected languages: {detected_info.get('languages', {})}")
        logger.debug(f"Detected topics: {detected_info.get('topics', {})}")
            
        matches = []
            
        # Search through all resources
        for file_path in self._content_cache:
            content = self._content_cache[file_path]
                    
                    # Calculate relevance score
            score = self._calculate_relevance_score(content, query, detected_info)
                    
            if score >= min_score:
                # Extract framework info from content or filename
                framework_info = self._extract_framework_info(content, file_path)
                        
                match = ResourceMatch(
                title=content.get('title', ''),
                language=content.get('language', ''),
                purpose=content.get('purpose', ''),
                code=content.get('code', ''),
                context=content.get('context', ''),
                source_file=file_path,
                topics=content.get('topics', []),
                framework_name=framework_info.get('name', ''),
                framework_version=framework_info.get('version', ''),
                difficulty_level=content.get('difficulty_level', ''),
                        score=score
                    )
                matches.append(match)
            
            logger.info(f"Found {len(matches)} matches above score {min_score}")
        return sorted(matches, key=lambda x: x.score, reverse=True)

    def _extract_framework_info(self, content: Dict, file_path: str) -> Dict[str, str]:
        """Extract framework name and version from content or filename."""
        framework_info = {'name': '', 'version': ''}
        
        # Try to extract from content first
        if 'framework' in content:
            framework_data = content['framework']
            if isinstance(framework_data, dict):
                framework_info['name'] = framework_data.get('name', '')
                framework_info['version'] = framework_data.get('version', '')
            elif isinstance(framework_data, str):
                framework_info['name'] = framework_data
        
        # If not found in content, try to extract from filename
        if not framework_info['name']:
            filename = os.path.basename(file_path).lower()
            common_frameworks = {
                'django': r'django[_-]?(\d+(?:\.\d+)*)?',
                'flask': r'flask[_-]?(\d+(?:\.\d+)*)?',
                'spring': r'spring[_-]?(\d+(?:\.\d+)*)?',
                'react': r'react[_-]?(\d+(?:\.\d+)*)?',
                'angular': r'angular[_-]?(\d+(?:\.\d+)*)?',
                'vue': r'vue[_-]?(\d+(?:\.\d+)*)?'
            }
            
            for framework, pattern in common_frameworks.items():
                match = re.search(pattern, filename)
                if match:
                    framework_info['name'] = framework
                    if match.group(1):
                        framework_info['version'] = match.group(1)
                    break
        
        return framework_info

    def generate_response(self, query: str, matches: List[ResourceMatch], year_level: str) -> Optional[str]:
        """Generate a response based on the search results."""
        if not matches:
            logger.info("No matches found for response generation")
            return None
        
        try:
            # Get the best match
            best_match = matches[0]
            
            # For simple program requests, provide a cleaner, tutorial-style response
            if "simple" in query.lower() and "program" in query.lower():
                response = f"""Let's create a simple {best_match.language.upper()} program that displays "Hello, world!" on the screen. This is the classic first program for almost every programmer.

Step 1: Understanding the Structure

A basic {best_match.language.upper()} program looks like this:

```{best_match.language}
{best_match.code}
```

This program introduces you to:
- Basic {best_match.language.upper()} program structure
- How to output text to the console
- The main() function that every {best_match.language.upper()} program needs"""
            else:
                # For other queries, show a more detailed but still clean response
                response = f"""Here's a {best_match.language.upper()} example for your question:

Purpose: {best_match.purpose}

```{best_match.language}
{best_match.code}
```

Key Points:
- Language: {best_match.language.upper()}
- Topics: {', '.join(best_match.topics)}
{f'- Context: {best_match.context}' if best_match.context else ''}"""

            logger.info("Generated response successfully")
            return response
            
        except Exception as e:
            logger.error(f"Error generating response: {str(e)}")
            return None

# Example usage
if __name__ == "__main__":
    retriever = SmartResourceRetriever()
    results = retriever.search_resources(
        "How to implement a binary search tree in Python?",
        "Year 2 Diploma"
    )
    if results:
        response = retriever.generate_response(
            "How to implement a binary search tree in Python?",
            results,
            "Year 2 Diploma"
        )
        print(response)
    else:
        print("No relevant resources found.")