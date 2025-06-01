# text_extractor.py
import os
import pdfplumber
import docx
import pptx
import re # For chunking
from text_cleaner import clean_documentation

# --- Chunking Parameters ---
CHUNK_SIZE = 700  # Characters per chunk (experiment with this)
CHUNK_OVERLAP = 100 # Overlap between chunks

def chunk_text(text: str, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP) -> list[str]:
    """Splits text into overlapping chunks."""
    if not text:
        return []
    # A simple way to split by paragraphs first, then by size if paragraphs are too long
    paragraphs = [p.strip() for p in text.splitlines() if p.strip()]
    chunks = []
    for para in paragraphs:
        if len(para) <= size:
            if para: # Ensure paragraph is not empty
                chunks.append(para)
        else:
            # If paragraph is too long, split it by fixed size with overlap
            start = 0
            while start < len(para):
                end = start + size
                chunks.append(para[start:end])
                start += (size - overlap)
    if not chunks and text: # Fallback if paragraph splitting yielded nothing but text exists
        start = 0
        while start < len(text):
            end = start + size
            chunks.append(text[start:end])
            start += (size - overlap)
    
    # Clean the chunks using our new cleaning functionality
    cleaned_chunks = clean_documentation([chunk for chunk in chunks if chunk])
    return cleaned_chunks

# --- Your existing extraction functions ---
def extract_text_from_pdf(file_path: str) -> str | None:
    try:
        with pdfplumber.open(file_path) as pdf:
            text_content = "".join([page.extract_text() + "\n" for page in pdf.pages if page.extract_text()])
            return text_content
    except Exception as e:
        print(f"Extractor: Error reading PDF {file_path}: {e}")
        return None

def extract_text_from_txt_or_code(file_path: str) -> str | None:
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"Extractor: Error reading text/code {file_path}: {e}")
        return None

def extract_text_from_docx(file_path: str) -> str | None:
    try:
        doc = docx.Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Extractor: Error reading DOCX {file_path}: {e}")
        return None

def extract_text_from_pptx(file_path: str) -> str | None:
    try:
        prs = pptx.Presentation(file_path)
        return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except Exception as e:
        print(f"Extractor: Error reading PPTX {file_path}: {e}")
        return None

def extract_text(file_path: str) -> str | None:
    filename = os.path.basename(file_path)
    file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
    if file_ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_ext in ["txt", "md", "cpp", "h", "js", "py", "gitignore", "log"]:
        return extract_text_from_txt_or_code(file_path)
    elif file_ext == "docx":
        return extract_text_from_docx(file_path)
    elif file_ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(file_path)
    else:
        return None