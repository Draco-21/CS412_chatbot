import os
import re
import pdfplumber  # For reading PDF files
import docx        # For reading DOCX files
import pptx        # For reading PPTX files

RESOURCES_BASE_FOLDER = "resources"
YEAR_FOLDER_MAP = {
    "Year 1 Certificate": "Year 1",
    "Year 2 Diploma": "Year 2",
    "Year 3 Degree": "Year 3",
    "Year 4 Postgraduate Diploma": "Year 4"
}

# --- Text Extraction Functions (No changes needed here) ---

def extract_text_from_pdf(file_path: str) -> str | None:
    """Extracts text content from a PDF file."""
    try:
        with pdfplumber.open(file_path) as pdf:
            text_content = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text + "\n"
            return text_content.lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading PDF file {file_path}: {e}")
        return None

def extract_text_from_txt_or_code(file_path: str) -> str | None:
    """Extracts text content from a .txt, .md, .cpp, .h, .js, .py file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading text/code file {file_path}: {e}")
        return None

def extract_text_from_docx(file_path: str) -> str | None:
    """Extracts text content from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text).lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading DOCX file {file_path}: {e}")
        return None

def extract_text_from_pptx(file_path: str) -> str | None:
    """Extracts text content from a PPTX file."""
    try:
        prs = pptx.Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text).lower()
    except Exception as e:
        print(f"ResourceRetriever: Error reading PPTX file {file_path}: {e}")
        return None

# --- NEW Helper Function to Search a Single Folder ---

def search_single_folder(folder_path: str, keywords: list[str]) -> str | None:
    """
    Searches recursively within a single folder path for files matching keywords.
    Returns the full content of the first matching file found.
    """
    if not os.path.isdir(folder_path):
        print(f"ResourceRetriever: Search path {folder_path} not found or not a directory.")
        return None

    try:
        for root, dirs, files in os.walk(folder_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                content = None
                file_ext = filename.lower().split('.')[-1] if '.' in filename else ''

                # Extract text based on file type
                if file_ext == "pdf":
                    content = extract_text_from_pdf(file_path)
                elif file_ext in ["txt", "md", "cpp", "h", "js", "py", "gitignore", "log"]:
                    content = extract_text_from_txt_or_code(file_path)
                elif file_ext == "docx":
                    content = extract_text_from_docx(file_path)
                elif file_ext in ["ppt", "pptx"]:
                    content = extract_text_from_pptx(file_path)

                if content:
                    # Check if ANY keyword is present
                    found_keywords_count = sum(1 for keyword in keywords if keyword in content)
                    if found_keywords_count > 0:
                        print(f"ResourceRetriever: Found {found_keywords_count} keyword(s) in '{file_path}'.")
                        return content # Return full content of the first match

    except Exception as e:
        print(f"ResourceRetriever: Error during search in {folder_path}: {e}")
        return None

    return None # No match found in this folder

# --- UPDATED Main Search Function ---

def find_answer_in_resources(user_message: str, year_level: str) -> str | None:
    """
    Searches for an answer within local resource files, checking the primary
    year level first, then expanding to other year levels.
    """
    print(f"ResourceRetriever: Main search for '{user_message}', primary year: {year_level}.")

    keywords = [word for word in re.split(r'\W+', user_message.lower()) if word]
    if not keywords:
        print("ResourceRetriever: No valid keywords in user message.")
        return None
    
    print(f"ResourceRetriever: Using keywords: {keywords}")

    # --- Determine Search Order ---
    primary_folder_name = YEAR_FOLDER_MAP.get(year_level)
    all_folder_names = list(YEAR_FOLDER_MAP.values())
    
    search_order_names = []
    
    # 1. Add primary folder first if it exists
    if primary_folder_name:
        search_order_names.append(primary_folder_name)
    else:
        print(f"ResourceRetriever: Warning - No folder mapping for {year_level}")

    # 2. Add all other folders (excluding the primary one if it was added)
    for name in all_folder_names:
        if name not in search_order_names:
            search_order_names.append(name)

    print(f"ResourceRetriever: Search order: {search_order_names}")

    # --- Execute Search in Order ---
    for folder_name in search_order_names:
        folder_path = os.path.join(RESOURCES_BASE_FOLDER, folder_name)
        print(f"ResourceRetriever: Now searching in: {folder_path}")
        
        result = search_single_folder(folder_path, keywords)
        
        if result:
            print(f"ResourceRetriever: Match found in {folder_path}.")
            return result # Return the first match we find

    # If we get here, no match was found in any folder
    print(f"ResourceRetriever: No match found in any folder for '{user_message}'.")
    return None

# --- Example Usage (remains the same for testing) ---
if __name__ == "__main__":
    # ... (Keep your example usage/dummy file creation if needed) ...
    print("Testing with 'C++ loops'")
    result = find_answer_in_resources("C++ loops", "Year 1 Certificate")
    if result:
        print("\n--- Search Result Found (Showing first 500 chars) ---")
        print(result[:500] + "...")
    else:
        print("\nNo local result found.")

    # Example: Test a query that might exist in Year 2 but searched from Year 1
    # Ensure you have a 'Year 2' folder with a 'java_syntax.txt' file for this test
    # year2_path = os.path.join(RESOURCES_BASE_FOLDER, "Year 2")
    # if not os.path.exists(year2_path):
    #     os.makedirs(year2_path)
    # with open(os.path.join(year2_path, "java_syntax.txt"), "w") as f:
    #      f.write("Java syntax is similar to C++ but has key differences.")

    # print("\nTesting with 'Java syntax' from Year 1")
    # result2 = find_answer_in_resources("Java syntax", "Year 1 Certificate")
    # if result2:
    #      print("\n--- Search Result Found ---")
    #      print(result2[:500] + "...")
    # else:
    #      print("\nNo local result found.")